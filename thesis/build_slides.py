"""Rebuilds the graduation-study presentation from the slide template.

    python thesis/build_slides.py

Reads thesis/template/deck_template.pptx, the pristine three-slide template
(title, contents, section divider), duplicates its divider slide out to the
full deck, and writes:

    thesis/presentation_slides.pptx

Every content slide follows the same shape: the template's title rule, two or
three short sentences underneath it, then the pictures. Detail that used to sit
on the slide lives in the SCRIPT table below, which becomes the speaker notes,
so the slides stay readable from the back of the room and nothing is lost.

Body text is one size throughout (BODY, 18 pt). Only the template's own title
slide and the picture captions sit outside that rule.

There are no section-divider slides. The slide that opens a section names it in
its own title ("Results: fragile objects"), and the contents page points at
that slide.

Figure sources are thesis/figures/: preprint_fig1 and 3-6 are used whole, and
the hardware plate (preprint_fig2) is cropped into its panels at build time, so
nothing here depends on generated files that are not in the repository. Where
no photograph exists yet, a labelled placeholder box marks what belongs there.

Requires: python-pptx, Pillow.
"""
import os
import re
import shutil
import tempfile
import zipfile
from xml.dom import minidom

from PIL import Image
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE

THESIS = os.path.dirname(os.path.abspath(__file__))
FIG = os.path.join(THESIS, "figures")
PHOTO = os.path.join(FIG, "photos")
TEMPLATE = os.path.join(THESIS, "template", "deck_template.pptx")
OUT = os.path.join(THESIS, "presentation_slides.pptx")

N_SLIDES = 17               # 3 template slides + 14 duplicated dividers
DIVIDER_SLIDE = "slide3.xml"   # the template slide every content slide starts from

# ---------------------------------------------------------------- palette ---
DEEP = RGBColor(0x2C, 0x0E, 0x63)   # template title purple
MID = RGBColor(0x3E, 0x14, 0x85)    # template accent purple
INK = RGBColor(0x1A, 0x0A, 0x3D)    # template body ink
LAV = RGBColor(0xD9, 0xC9, 0xF5)    # template light lavender
CARD = RGBColor(0xF4, 0xF0, 0xFC)   # card fill, derived from LAV
GREY = RGBColor(0xF2, 0xF2, 0xF4)
MUTE = RGBColor(0x6B, 0x66, 0x7B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

# ------------------------------------------------------------------ type ---
BODY = 18       # every sentence, label and table cell on every content slide
CAP = 11        # picture captions and placeholder labels only

# ------------------------------------------------------- template geometry ---
CL, CW = 900450, 7343100            # content band, from the template body box
CT, CB = 1405800, 4690000
GAP = 228600
COLW = (CW - GAP) // 2
CR = CL + CW
LINE18 = 312000                     # one 18 pt Calibri line at 112%, in EMU
CPL = 70                            # 18 pt characters that fit across CW

# =============================================================== script ======
# The slides carry the headline; these carry everything else, as speaker notes.
SCRIPT = {}

SCRIPT[1] = """
Good afternoon. I am Adriel Santoso, from the Intelligent Control Systems Lab.

My graduation study asks a simple question: if you are driving a robot hand from
a distance, does giving your own fingers a sense of touch make you better at
picking up fragile things, and does it matter which kind of buzzer you use?
"""

SCRIPT[2] = """
Here is where we are going: the problem, the system, how the experiment was run,
what came out of it, and what it does and does not prove.
"""

SCRIPT[3] = """
Let me start with why this is hard.

Before anything else, the five words in my title, because most of you are not
from robotics.

Teleoperation is driving a machine from somewhere else: a surgeon at a console,
a technician outside a glovebox. Tactile means touch: information that only
exists at the moment of physical contact. Haptic is the return trip: a sensation
delivered back to your own skin.

A sensor turns contact into a number; an actuator turns a number back into
something you can feel. This study puts a sensor on the robot's fingers and an
actuator on yours, and asks whether closing that loop changes how you grasp.
"""

SCRIPT[4] = """
Think about picking up an egg. You do not calculate anything: receptors in your
skin report how hard you are squeezing within milliseconds, and you stop short of
cracking the shell.

Now take that away. In teleoperation you see the object but feel nothing, so grip
force has to be inferred from how it looks. That works right up until the object
does not visibly change before it fails. Too little force and it slips; too much
and you get the picture on the right.

And the gripper cannot rescue you. The Robotiq 2F-85 estimates force from one
aggregate motor current, with no contact distribution at all, and on my unit that
register reads zero milliamps regardless of contact.
"""

SCRIPT[5] = """
So the study asks two questions.

First: does touch feedback change how people grasp fragile and deformable
objects, compared with vision alone?

Second, and this is asked less often: does the actuator type matter? Most
published work picks one haptic device and compares it against a no-feedback
baseline. Here the whole sensing and control pipeline is held fixed, and the only
thing that changes is the hardware on the skin: a vibration motor at the finger
joints, or an electromagnetic pin at the fingertip. Three conditions: vision
only, LRA, EM.
"""

SCRIPT[6] = """
Here is the system. One closed loop, with a human inside it.

Going down: a camera watches the operator's hand, MediaPipe gives me the pinch
distance in pixels, and that maps to a gripper opening sent over Modbus at 25
hertz. The gripper mirrors your pinch. There is no joystick.

Coming back up: each jaw carries a soft tactile sensor. Fifteen times a second I
read how deeply the gel is dented, turn it into one number between zero and one,
and stream it to the wrist.

The loop closes through the person, and there is no force controller anywhere in
it. Two guards: closing stops at one millimetre of indentation, and the actuators
cut out after 200 milliseconds of silence.
"""

SCRIPT[7] = """
The hardware. The operator wears a wrist unit carrying the driver board and five
independent actuator channels, and a plain USB camera does the hand tracking,
which is two-dimensional and uncalibrated.

On the robot side: the Robotiq 2F-85 gripper with a plastic egg in the fixture,
and a 9DTact vision-based tactile sensor on each jaw, about the size of a sugar
cube.
"""

SCRIPT[8] = """
This is where the engineering actually lives: turning a squishy gel into a
number.

The sensor is a camera looking at the back of a soft gel pad. Every frame it
reconstructs a height map; subtract the map recorded with nothing touching it and
you have the deformation field. Two quantities come out of that one field.

The first is the summed deformation over every pixel in contact. That is my
grip-force proxy: monotonic in normal force, but uncalibrated, so I report
arbitrary units rather than newtons.

The second is the 99th-percentile depth, and that is what the operator feels:
divide by a saturation depth, clip to zero-to-one. Two millimetres for the egg,
six tenths for foam, which barely dents the gel.
"""

SCRIPT[9] = """
Two actuators render that number.

The LRA is a linear resonant actuator, the vibration motor in your phone. A
bipolar carrier at 200 hertz, with intensity setting the envelope, so it buzzes
harder as you squeeze. It sits on the proximal joints and draws power throughout.

The EM is a bistable pin on an H-bridge. A four-millisecond pulse throws the pin
into the skin and it latches, so it holds at zero power; intensity sets the gap
between bursts. It sits on the fingertips.

Note what I just said: the two conditions differ in technology and in placement.
Hold on to that, because it comes back in the discussion.
"""

SCRIPT[10] = """
The experiment.

Twenty-two participants, each doing all three conditions in the fixed order:
vision only, then LRA, then EM. Two object classes, five grasps each, so 660 trials
in total.

Fragile is a hollow plastic egg that pops apart at the seam: it fails suddenly,
like the real thing, but I can reassemble it. Deformable is a foam cube, which
yields gradually and never breaks. Afterwards each participant rated the three
conditions on four items and picked a favourite.
"""

SCRIPT[11] = """
Four objective measures. Peak deformation is how hard they squeezed at the worst
moment.

Excess deformation is the interesting one. Once the gel is fully dented the depth
signal plateaus, and squeezing harder past that point changes nothing the operator
can feel. That wasted, potentially destructive effort is what this measures.

Grip adjustments counts force reversals after the plateau: somebody hunting for a
grip instead of holding steady. And survival: did the egg come out intact.

Friedman across the three conditions, then Wilcoxon with Holm correction. Medians,
not means, because one crushed egg should not move the average.
"""

SCRIPT[12] = """
Results. Fragile objects, median participant, and for the first three rows lower
is better.

Excess deformation, the wasted squeezing, falls from about 1530 units to 170 with
the LRA, roughly ninefold, at p equals 0.019. Grip adjustments drop from seven to
two with the EM, p equals 0.021. The two actuators help in different ways: the LRA
stops you squeezing, the EM stops you fidgeting.

And the outcome people care about: survival rises from 62 percent to 81 with the
LRA and 78 with the EM, Friedman p equals 0.004, both beating the baseline after
correction. LRA versus EM is 0.69.

Do not read the table aloud. Point at three numbers: 1530 to 170, 7 to 2, 62 to
81.
"""

SCRIPT[13] = """
The same result as distributions rather than medians.

Two honest negatives live here. Peak dent depth barely moves, because my
one-millimetre safety cutoff caps it before the conditions can differ. And
deformable objects showed no reliable change on anything, because foam gives you no
sharp failure to avoid, so there is nothing for the feedback to buy you.
"""

SCRIPT[14] = """
Subjectively it is not close. On force perception, grasp confidence and contact
detection, both haptic conditions beat vision only at p below 0.001, surviving
Holm correction.

But no single question separated the LRA from the EM: every pairwise comparison
between them sits above 0.2. And yet, asked to pick a favourite, 17 of 22 chose
the LRA, 4 the EM and 1 vision only. People had a clear preference they could not
articulate on any of my rating items.
"""

SCRIPT[15] = """
So what does this mean. What I am comfortable claiming: tactile feedback reduces
the squeezing the operator cannot feel, reduces regrip hunting, and raises
survival on fragile objects.

What I am not comfortable claiming is that the LRA is the better actuator.
Technology and placement are bundled: I never tested a fingertip LRA or a
joint-mounted EM. The fingertip EM also sits on the landmarks the hand tracker
uses, so the preference may be smoother control rather than better sensation. And
the condition order was fixed, so practice and fatigue are inseparable from
condition.

On measurement: force is in arbitrary units, and latency still needs a bench
measurement. Every test here is rank-based, so a scaling error cannot move a
p-value.

Expect questions on the uncalibrated force and the fixed order. Answer them from
this slide.
"""

SCRIPT[16] = """
To close. Three things.

Touch feedback made remote grasping of fragile objects measurably safer: about
nine times less unfelt squeezing, and survival up from 62 to around 80 percent.

It is what operators want: 21 of 22 preferred a haptic condition over vision
alone.

And which actuator is better is still open; answering it needs a study that
separates technology from placement.
"""

SCRIPT[17] = """
Where this goes next.

The setup is a laboratory one, but the problem is not. Anywhere a person drives a
gripper they cannot feel, and the payload is fragile, the same failure is waiting:
keyhole surgery, glovebox and hot-cell handling, disaster response, orbital
servicing. None of those need a new robot. They need a better signal on the skin.

My own next steps are three: spatial feedback instead of one number per finger,
force calibrated into newtons so the numbers transfer between rigs, and the
crossed study that separates actuator technology from where it sits on the hand.

Thank you, and happy to take questions.
"""


# =============================================================== helpers =====
def tidy(s):
    """Collapses the triple-quoted script blocks into paragraphs."""
    paras = [" ".join(p.split()) for p in s.strip().split("\n\n")]
    return "\n\n".join(p for p in paras if p)


def title_shape(slide):
    return slide.shapes[0]


def body_shape(slide):
    return slide.shapes[4]


def set_title(slide, text):
    """Replaces the title run text, keeping the template's run formatting."""
    tf = title_shape(slide).text_frame
    p = tf.paragraphs[0]
    runs = p.runs
    runs[0].text = text
    for r in runs[1:]:
        r._r.getparent().remove(r._r)


def drop_body(slide):
    sp = body_shape(slide)._element
    sp.getparent().remove(sp)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def textbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Emu(int(x)), Emu(int(y)), Emu(int(w)),
                                  Emu(int(h)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    return tb


def para(tf, first, runs, size=BODY, color=INK, bold=False, space_before=0,
         space_after=6, align=PP_ALIGN.LEFT, line=112, hang=0):
    """Adds a paragraph made of (text, bold, color, size) run tuples."""
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = line / 100.0
    if hang:
        pPr = p._p.get_or_add_pPr()
        pPr.set("marL", str(hang))
        pPr.set("indent", str(-hang))
    if isinstance(runs, str):
        runs = [(runs, bold, color, size)]
    for text, b, c, s in runs:
        r = p.add_run()
        r.text = text
        f = r.font
        f.name = FONT
        f.size = Pt(s)
        f.bold = b
        f.color.rgb = c
    return p


def block(slide, x, y, w, h, items, anchor=MSO_ANCHOR.TOP):
    """items: list of dicts -> {runs|text, size, color, bold, space_*, align}."""
    tf = textbox(slide, x, y, w, h, anchor).text_frame
    for i, it in enumerate(items):
        it = dict(it)
        runs = it.pop("runs", None)
        if runs is None:
            runs = it.pop("text")
        else:
            it.pop("text", None)
        para(tf, i == 0, runs, **it)
    return tf


def lead(slide, sentences, y=CT, x=CL, w=CW, size=BODY):
    """The two-or-three sentence opener every content slide starts with.

    Returns the y coordinate just below the text, so the pictures can be laid
    out underneath it the way the reference deck does.
    """
    if isinstance(sentences, str):
        sentences = [sentences]
    h = 0
    items = []
    for i, s in enumerate(sentences):
        items.append(dict(text=s, size=size, color=INK, line=112,
                          space_after=6 if i < len(sentences) - 1 else 0))
        # CPL characters fit on one line across the full content band
        per_line = max(8, CPL * w // CW)
        h += LINE18 * max(1, -(-len(s) // per_line)) + 76200
    block(slide, x, y, w, h, items)
    return y + h + 100000


def card(slide, x, y, w, h, fill=CARD, line=LAV, radius=0.08):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Emu(int(x)), Emu(int(y)), Emu(int(w)),
                                Emu(int(h)))
    sh.adjustments[0] = radius
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ""
    return sh


def pic_fit(slide, path, x, y, boxw, boxh, center=True):
    """Fits an image inside a box, centred, and returns (x, y, w, h)."""
    iw, ih = Image.open(path).size
    scale = min(boxw / iw, boxh / ih)
    w, h = int(iw * scale), int(ih * scale)
    px = x + (boxw - w) // 2 if center else x
    py = y + (boxh - h) // 2 if center else y
    slide.shapes.add_picture(path, Emu(int(px)), Emu(int(py)), Emu(int(w)),
                             Emu(int(h)))
    return px, py, w, h


def placeholder(slide, x, y, w, h, label):
    """A dashed box standing in for a photograph that does not exist yet.

    The label says what belongs there, so the box can be swapped for a real
    image without having to re-read the script.
    """
    sh = card(slide, x, y, w, h, fill=GREY, line=RGBColor(0xBD, 0xB6, 0xCC))
    sh.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    block(slide, x + 137160, y, w - 274320, h, [
        dict(text="IMAGE PLACEHOLDER", size=CAP, bold=True, color=MUTE,
             align=PP_ALIGN.CENTER, space_after=4, line=105),
        dict(text=label, size=CAP, color=MUTE, align=PP_ALIGN.CENTER,
             space_after=0, line=110)], anchor=MSO_ANCHOR.MIDDLE)
    return x, y, w, h


def caption(slide, x, y, w, text, align=PP_ALIGN.CENTER):
    block(slide, x, y, w, 240000, [
        dict(text=text, size=CAP, color=MUTE, align=align, space_after=0,
             line=110)])


def plate(slide, y, entries, boxh, gap=GAP, name_size=BODY):
    """A row of pictures with a bold name and a small caption under each.

    entries: (path_or_None, name, caption_text, placeholder_label) tuples.
    """
    n = len(entries)
    w = (CW - (n - 1) * gap) // n
    for i, (path, name, cap, ph_label) in enumerate(entries):
        x = CL + i * (w + gap)
        if path is None:
            placeholder(slide, x, y, w, boxh, ph_label)
        else:
            pic_fit(slide, path, x, y, w, boxh)
        ty = y + boxh + 60000
        if name:
            block(slide, x, ty, w, 300000, [
                dict(text=name, size=name_size, bold=True, color=DEEP,
                     align=PP_ALIGN.CENTER, space_after=0)])
            ty += 300000
        if cap:
            caption(slide, x, ty, w, cap)
    return w


def arrow(slide, x, y, w, h):
    sh = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x)), Emu(int(y)),
                                Emu(int(w)), Emu(int(h)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = LAV
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def connect(slide, x1, y1, x2, y2):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(int(x1)),
                                    Emu(int(y1)), Emu(int(x2)), Emu(int(y2)))
    ln.line.color.rgb = LAV
    ln.line.width = Pt(1.5)
    ln.shadow.inherit = False
    for ref in ln._element.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}"
            "effectRef"):
        ref.set("idx", "0")
    return ln


# ================================================================ build ======
# ======================================================= deck preparation ====
# The template ships three slides; every content slide in the finished deck is
# a copy of its section-divider slide, so the title rule, the paired logos and
# the slide-number placeholder stay identical across the deck. Duplicating a
# slide means copying the part, its rels, the content-type override and the
# entry in <p:sldIdLst>, and python-pptx cannot do this, so it is done here on the
# unpacked package. The notesSlide relationship is deliberately dropped from
# the copies: each new slide gets its own notes when notes are written to it.
NS_CT = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"
REL_SLIDE = ("http://schemas.openxmlformats.org/officeDocument/2006/"
             "relationships/slide")
REL_NOTES = ("http://schemas.openxmlformats.org/officeDocument/2006/"
             "relationships/notesSlide")


def _next_free(existing, pattern, start=1):
    n = start
    while pattern % n in existing:
        n += 1
    return n


def duplicate_slide(root, source):
    """Copies ppt/slides/<source> into the package, returning the new name."""
    slide_dir = os.path.join(root, "ppt", "slides")
    taken = set(os.listdir(slide_dir))
    new_name = "slide%d.xml" % _next_free(taken, "slide%d.xml")
    shutil.copyfile(os.path.join(slide_dir, source),
                    os.path.join(slide_dir, new_name))

    # slide rels: copy, minus the notes slide (each copy gets fresh notes)
    rels_dir = os.path.join(slide_dir, "_rels")
    src_rels = os.path.join(rels_dir, source + ".rels")
    if os.path.exists(src_rels):
        doc = minidom.parse(src_rels)
        for rel in list(doc.getElementsByTagNameNS(NS_REL, "Relationship")):
            if rel.getAttribute("Type") == REL_NOTES:
                rel.parentNode.removeChild(rel)
        with open(os.path.join(rels_dir, new_name + ".rels"), "wb") as fh:
            fh.write(doc.toxml("UTF-8"))

    # content-type override (attribute order in the template is not assumed)
    ct_path = os.path.join(root, "[Content_Types].xml")
    ct = open(ct_path, encoding="utf-8").read()
    override = ('<Override PartName="/ppt/slides/%s" ContentType='
                '"application/vnd.openxmlformats-officedocument.'
                'presentationml.slide+xml"/>') % new_name
    ct = ct.replace("</Types>", override + "</Types>")
    open(ct_path, "w", encoding="utf-8").write(ct)

    # presentation rels + <p:sldIdLst> entry
    pres_rels_path = os.path.join(root, "ppt", "_rels", "presentation.xml.rels")
    pres_rels = open(pres_rels_path, encoding="utf-8").read()
    rid = "rId%d" % _next_free(set(re.findall(r'Id="(rId\d+)"', pres_rels)),
                               "rId%d")
    pres_rels = pres_rels.replace(
        "</Relationships>",
        '<Relationship Id="%s" Type="%s" Target="slides/%s"/></Relationships>'
        % (rid, REL_SLIDE, new_name))
    open(pres_rels_path, "w", encoding="utf-8").write(pres_rels)

    pres_path = os.path.join(root, "ppt", "presentation.xml")
    pres = open(pres_path, encoding="utf-8").read()
    sld_id = max(int(i) for i in re.findall(r'<p:sldId id="(\d+)"', pres)) + 1
    pres = pres.replace(
        "</p:sldIdLst>",
        '<p:sldId id="%d" r:id="%s"/></p:sldIdLst>' % (sld_id, rid))
    open(pres_path, "w", encoding="utf-8").write(pres)
    return new_name


def build_base(workdir):
    """Unpacks the template, grows it to N_SLIDES, repacks it."""
    root = os.path.join(workdir, "package")
    with zipfile.ZipFile(TEMPLATE) as zf:
        zf.extractall(root)
    n_existing = len(os.listdir(os.path.join(root, "ppt", "slides"))) - 1
    for _ in range(N_SLIDES - n_existing):
        duplicate_slide(root, DIVIDER_SLIDE)
    base = os.path.join(workdir, "base.pptx")
    with zipfile.ZipFile(base, "w", zipfile.ZIP_DEFLATED) as zf:
        for folder, _, files in os.walk(root):
            for name in files:
                full = os.path.join(folder, name)
                zf.write(full, os.path.relpath(full, root))
    return base


def crop_panels(workdir):
    """Cuts the hardware plate (preprint_fig2) into the panels the deck uses.

    Slides 7 and 9 show the apparatus and the two actuators separately, so
    the plate's (a)/(b) sub-labels are cropped away here rather than being
    edited into the figure; preprint_fig2.png stays the preprint's copy.
    """
    plate_img = Image.open(os.path.join(FIG, "preprint_fig2.png"))
    w, _ = plate_img.size
    out = {}
    panels = {
        "apparatus": plate_img.crop((0, 88, w, 900)),      # (a), label removed
        "lra": plate_img.crop((95, 1015, 625, 1712)),      # (b) left half
        "em": plate_img.crop((655, 1015, 1225, 1712)),     # (b) right half
    }
    for name, img in panels.items():
        path = os.path.join(workdir, name + ".png")
        img.save(path)
        out[name] = path
    return out


WORK = tempfile.mkdtemp(prefix="deckbuild-")
try:
    PANEL = crop_panels(WORK)
    prs = Presentation(build_base(WORK))
    S = list(prs.slides)


    def s(n):
        return S[n - 1]


    # --- 1. title slide: untouched template ------------------------------------
    notes(s(1), tidy(SCRIPT[1]))

    # --- 2. contents ------------------------------------------------------------
    sl = s(2)
    drop_body(sl)
    entries = [("Introduction & Background", 3), ("System", 6), ("Method", 10),
               ("Results", 12), ("Discussion", 15), ("Conclusion", 16)]
    y = 1460000
    row_h = 500000
    for label, page in entries:
        block(sl, CL, y, CW - 700000, row_h, [
            dict(text=label, size=BODY, color=INK, space_after=0)])
        ln = sl.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(CL + 3100000), Emu(y + 150000),
            Emu(CR - 620000), Emu(y + 150000))
        ln.line.color.rgb = LAV
        ln.line.width = Pt(1)
        ln.line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT
        ln.shadow.inherit = False
        for ref in ln._element.findall(
                ".//{http://schemas.openxmlformats.org/drawingml/2006/main}"
                "effectRef"):
            ref.set("idx", "0")
        block(sl, CR - 500000, y, 500000, row_h, [
            dict(text=str(page), size=BODY, color=MID, bold=True,
                 align=PP_ALIGN.RIGHT, space_after=0)])
        y += row_h
    notes(sl, tidy(SCRIPT[2]))

    # --- 3. the title, in plain words -------------------------------------------
    # Opens the introduction, so it carries the section name in its own title.
    # A branch off the study title: each word in it gets one plain-English line.
    sl = s(3)
    set_title(sl, "Introduction: the words in the title")
    drop_body(sl)
    rootw, rgap = 2050000, 400000
    lx = CL + rootw + rgap
    lw = CR - lx
    leaf_h, leaf_gap = 600000, 60000
    terms = [
        ("Teleoperation", "driving a machine from somewhere else"),
        ("Tactile", "to do with touch, physical contact"),
        ("Haptic", "a sensation delivered back to your skin"),
        ("Sensor", "turns contact into a number"),
        ("Actuator", "turns a number into something you feel"),
    ]
    total = len(terms) * leaf_h + (len(terms) - 1) * leaf_gap
    ly0 = CT + (CB - CT - total) // 2
    rooth = 1250000
    rooty = CT + (CB - CT - rooth) // 2
    card(sl, CL, rooty, rootw, rooth, fill=LAV, line=MID)
    block(sl, CL + 130000, rooty, rootw - 260000, rooth, [
        dict(text="Tactile-feedback teleoperation", size=BODY, bold=True,
             color=DEEP, align=PP_ALIGN.CENTER, space_after=0, line=112)],
        anchor=MSO_ANCHOR.MIDDLE)
    for i, (term, gloss) in enumerate(terms):
        ly = ly0 + i * (leaf_h + leaf_gap)
        connect(sl, CL + rootw, rooty + rooth // 2, lx, ly + leaf_h // 2)
        card(sl, lx, ly, lw, leaf_h, fill=CARD, line=LAV)
        block(sl, lx + 150000, ly, lw - 300000, leaf_h, [
            dict(runs=[(term + ": ", True, MID, BODY),
                       (gloss, False, INK, BODY)], space_after=0, line=110)],
            anchor=MSO_ANCHOR.MIDDLE)
    notes(sl, tidy(SCRIPT[3]))

    # --- 4. the missing sense ---------------------------------------------------
    sl = s(4)
    set_title(sl, "Teleoperation has no sense of touch")
    drop_body(sl)
    y = lead(sl, [
        "Your skin reports squeeze force in milliseconds, so you stop in time.",
        "A remote operator has only video, and must guess the grip force.",
        "Too little and the object slips; too much and it breaks."])
    plate(sl, y, [
        (os.path.join(PHOTO, "object_fragile.png"), "",
         "The fragile object: a hollow shell", None),
        (os.path.join(PHOTO, "object_broken.png"), "",
         "The same shell, broken at the seam", None),
        (None, "", "What the operator actually works from",
         "PHOTO: the operator's console, screen and hand only"),
    ], CB - y - 320000)
    notes(sl, tidy(SCRIPT[4]))

    # --- 5. two questions -------------------------------------------------------
    sl = s(5)
    set_title(sl, "Two questions, one fixed pipeline")
    drop_body(sl)
    y = lead(sl, [
        "Does touch feedback change how people grasp fragile objects?",
        "And does it matter which actuator delivers that touch?"])
    plate(sl, y, [
        (os.path.join(PHOTO, "camera.png"), "Vision only",
         "baseline: the screen and nothing else", None),
        (os.path.join(PHOTO, "lra_module.png"), "LRA",
         "vibration motor at the finger joints", None),
        (os.path.join(PHOTO, "em_cap.png"), "EM",
         "magnetic pin tapping the fingertip", None),
    ], CB - y - 620000)
    notes(sl, tidy(SCRIPT[5]))

    # --- 6. control loop --------------------------------------------------------
    # Opens the system section, so it carries the section name in its title.
    sl = s(6)
    set_title(sl, "System: one closed loop, human included")
    drop_body(sl)
    y = lead(sl, [
        "A camera reads the operator's pinch; the gripper mirrors it at 25 Hz.",
        "The gel sensors send depth back at 15 Hz, and the person eases off."])
    plate(sl, y, [
        (os.path.join(FIG, "preprint_fig1.png"), "", "The loop, as a diagram",
         None),
        (os.path.join(PHOTO, "lra_hand.png"), "", "Operator side: actuators on "
         "the hand", None),
        (os.path.join(PHOTO, "gripper.png"), "", "Robot side: gripper and gel "
         "sensors", None),
    ], CB - y - 320000)
    notes(sl, tidy(SCRIPT[6]))

    # --- 7. apparatus -----------------------------------------------------------
    sl = s(7)
    set_title(sl, "The apparatus")
    drop_body(sl)
    y = lead(sl, [
        "The operator wears the driver board and is tracked by a USB camera.",
        "The robot is a Robotiq 2F-85 with a 9DTact gel sensor on each jaw."])
    plate(sl, y, [
        (os.path.join(PHOTO, "wearable_hand.png"), "",
         "Wearable: ESP32-C6, five actuator channels", None),
        (PANEL["apparatus"], "", "The full bench, as run", None),
        (os.path.join(PHOTO, "touch_sensor.png"), "",
         "One gel sensor, about a sugar cube", None),
    ], CB - y - 320000)
    notes(sl, tidy(SCRIPT[7]))

    # --- 8. gel to intensity ----------------------------------------------------
    sl = s(8)
    set_title(sl, "From gel dent to felt intensity")
    drop_body(sl)
    y = lead(sl, [
        "A camera behind the gel rebuilds its surface every frame.",
        "Subtracting the untouched baseline leaves the deformation.",
        "Its sum is the force proxy; its depth is what the operator feels."])
    aw = 320000
    boxw = (CW - 2 * aw) // 3
    boxh = min(1850000, CB - y - 460000)
    labels = [
        ("PHOTO: raw gel image under the sensor's illumination",
         "Camera image of the gel"),
        ("RENDER: reconstructed height map h(x, y)", "Height map"),
        ("RENDER: deformation field d = h − h₀, contact masked",
         "Deformation field"),
    ]
    for i, (ph_label, cap) in enumerate(labels):
        x = CL + i * (boxw + aw)
        placeholder(sl, x, y, boxw, boxh, ph_label)
        caption(sl, x, y + boxh + 60000, boxw, cap)
        if i < 2:
            arrow(sl, x + boxw + 60000, y + boxh // 2 - 90000, aw - 120000,
                  180000)
    notes(sl, tidy(SCRIPT[8]))

    # --- 9. two actuators -------------------------------------------------------
    sl = s(9)
    set_title(sl, "Two ways to render that number")
    drop_body(sl)
    y = lead(sl, [
        "The LRA buzzes at the finger joints; the EM taps the fingertip.",
        "They differ in technology and in placement, so remember that."])
    plate(sl, y, [
        (PANEL["lra"], "LRA", "linear resonant actuator, proximal joints", None),
        (PANEL["em"], "EM", "bistable magnetic pin, fingertips", None),
    ], CB - y - 620000)
    notes(sl, tidy(SCRIPT[9]))

    # --- 10. protocol -----------------------------------------------------------
    # Opens the method section, so it carries the section name in its title.
    sl = s(10)
    set_title(sl, "Method: how the study ran")
    drop_body(sl)
    y = lead(sl, [
        "Everyone ran all three conditions, in the fixed order vision → LRA → EM.",
        "660 trials: 5 grasps × 2 objects × 3 conditions × 22 people.",
        "Then they rated each condition and picked a favourite."])
    boxh = CB - y - 300000
    w = (CW - GAP) // 2
    pic_fit(sl, os.path.join(FIG, "preprint_fig3.png"), CL, y, w, boxh)
    caption(sl, CL, y + boxh + 60000, w, "Session flow and condition order")
    pic_fit(sl, os.path.join(FIG, "preprint_fig4.png"), CL + w + GAP, y, w,
            boxh)
    caption(sl, CL + w + GAP, y + boxh + 60000, w,
            "Fragile: a shell that pops apart.  Deformable: foam that never breaks.")
    notes(sl, tidy(SCRIPT[10]))

    # --- 11. metrics ------------------------------------------------------------
    sl = s(11)
    set_title(sl, "What was measured")
    drop_body(sl)
    y = lead(sl, [
        "Four measures per trial, all of them rank-tested."])
    mw = 5000000
    metrics = [
        ("Peak deformation", "the hardest squeeze"),
        ("Excess deformation", "force you cannot feel"),
        ("Grip adjustments", "reversals after that plateau"),
        ("Survival", "did the object come out intact"),
    ]
    mh = (CB - y - 3 * 90000) // 4
    for i, (name, gloss) in enumerate(metrics):
        my = y + i * (mh + 90000)
        card(sl, CL, my, mw, mh, fill=CARD, line=LAV)
        block(sl, CL + 150000, my, mw - 300000, mh, [
            dict(runs=[(name + ": ", True, MID, BODY),
                       (gloss, False, INK, BODY)], space_after=0, line=110)],
            anchor=MSO_ANCHOR.MIDDLE)
    px = CL + mw + GAP
    pw = CR - px
    placeholder(sl, px, y, pw, CB - y - 300000,
                "PLOT: one force trace, with the plateau and the post-plateau "
                "reversals marked")
    caption(sl, px, CB - 230000, pw, "Where the metrics come from")
    notes(sl, tidy(SCRIPT[11]))

    # --- 12. objective results --------------------------------------------------
    # Opens the results section, so it carries the section name in its title.
    sl = s(12)
    set_title(sl, "Results: fragile objects")
    drop_body(sl)
    y = lead(sl, [
        "Median participant, three conditions; on the first three rows, lower is "
        "better."])
    rows = [
        ("Median per participant", "Vision", "LRA", "EM"),
        ("Peak deformation volume (a.u.)", "20100", "12300", "16000"),
        ("Excess deformation (a.u.)", "1530", "170", "260"),
        ("Grip adjustments (#)", "7", "5", "2"),
        ("Objects surviving (%)", "62", "81", "78"),
    ]
    tbl_h = min(2200000, CB - y - 300000)
    gt = sl.shapes.add_table(len(rows), 4, Emu(CL), Emu(int(y)), Emu(CW),
                             Emu(int(tbl_h))).table
    gt.columns[0].width = Emu(3400000)
    for c in range(1, 4):
        gt.columns[c].width = Emu((CW - 3400000) // 3)
    for r in range(len(rows)):
        gt.rows[r].height = Emu(int(tbl_h // len(rows)))
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.margin_left = Emu(91440)
            cell.margin_right = Emu(91440)
            cell.margin_top = Emu(18288)
            cell.margin_bottom = Emu(18288)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r == 0:
                cell.fill.fore_color.rgb = DEEP
            elif r == len(rows) - 1:
                cell.fill.fore_color.rgb = LAV
            else:
                cell.fill.fore_color.rgb = WHITE if r % 2 else CARD
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run()
            run.text = val
            run.font.name = FONT
            run.font.size = Pt(BODY)
            run.font.bold = (r == 0) or (r == len(rows) - 1) or (c == 0)
            run.font.color.rgb = WHITE if r == 0 else INK
    caption(sl, CL, y + tbl_h + 80000, CW,
            "Excess deformation vision→LRA p = 0.019  ·  grip adjustments "
            "vision→EM p = 0.021  ·  survival Friedman p = 0.004",
            align=PP_ALIGN.LEFT)
    notes(sl, tidy(SCRIPT[12]))

    # --- 13. the same result, as distributions ----------------------------------
    sl = s(13)
    set_title(sl, "The same result, spread out")
    drop_body(sl)
    y = lead(sl, [
        "Both actuators cut the wasted squeezing, and the LRA cuts it hardest.",
        "Two honest negatives: the safety cutoff caps peak depth, and foam "
        "moved on nothing."])
    _, iy, _, ih = pic_fit(sl, os.path.join(FIG, "preprint_fig5.png"),
                           CL, y, CW, CB - y - 260000)
    caption(sl, CL, iy + ih + 70000, CW,
            "Per-condition distributions on the fragile object")
    notes(sl, tidy(SCRIPT[13]))

    # --- 14. subjective results -------------------------------------------------
    sl = s(14)
    set_title(sl, "What the participants said")
    drop_body(sl)
    y = lead(sl, [
        "Both haptic conditions beat vision alone on every item (p < 0.001).",
        "No item separated the two, yet 17 of 22 picked the LRA."])
    _, iy, _, ih = pic_fit(sl, os.path.join(FIG, "preprint_fig6.png"),
                           CL, y, CW, CB - y - 260000)
    caption(sl, CL, iy + ih + 70000, CW,
            "Likert ratings by condition, and the forced-choice preference")
    notes(sl, tidy(SCRIPT[14]))

    # --- 15. limitations --------------------------------------------------------
    # Opens the discussion section, so it carries the section name in its title.
    sl = s(15)
    set_title(sl, "Discussion: claims and confounds")
    drop_body(sl)
    y = lead(sl, [
        "One column I will defend; the other I will not, at least not yet."])
    ch = CB - y - 60000
    card(sl, CL, y, COLW, ch, fill=CARD, line=LAV)
    block(sl, CL + 150000, y + 140000, COLW - 300000, ch - 280000, [
        dict(text="Supported", size=BODY, bold=True, color=DEEP, space_after=12),
        dict(text="Unfelt squeezing falls ninefold.", space_after=12, line=110),
        dict(text="Survival rises 62% → 80%.", space_after=12, line=110),
        dict(text="Operators notice it and want it.", space_after=0, line=110),
    ])
    x2 = CL + COLW + GAP
    card(sl, x2, y, COLW, ch, fill=GREY, line=RGBColor(0xDD, 0xDD, 0xE2))
    block(sl, x2 + 150000, y + 140000, COLW - 300000, ch - 280000, [
        dict(text="Not yet supported", size=BODY, bold=True, color=DEEP,
             space_after=12),
        dict(text="That the LRA is truly the better one.", space_after=12,
             line=110),
        dict(text="That any of it holds for foam.", space_after=12, line=110),
        dict(text="That the fixed order did not matter.", space_after=0,
             line=110),
    ])
    notes(sl, tidy(SCRIPT[15]))

    # --- 16. conclusion ---------------------------------------------------------
    # Opens the conclusion section, so it carries the section name in its title.
    sl = s(16)
    set_title(sl, "Conclusion: takeaways")
    drop_body(sl)
    takeaways = [
        ("62% → 80%", "Fragile objects survived remote grasping far more often "
         "with touch feedback."),
        ("21 of 22", "Participants preferred a haptic condition over vision "
         "alone."),
        ("Still open", "Which actuator is better: placement has to be separated "
         "from technology."),
    ]
    th = (CB - CT - 2 * 140000) // 3
    for i, (big, text) in enumerate(takeaways):
        y = CT + i * (th + 140000)
        card(sl, CL, y, CW, th, fill=CARD if i < 2 else WHITE, line=LAV)
        block(sl, CL + 150000, y, 1900000, th,
              [dict(text=big, size=BODY, bold=True, color=DEEP, space_after=0)],
              anchor=MSO_ANCHOR.MIDDLE)
        block(sl, CL + 2150000, y, CW - 2300000, th,
              [dict(text=text, size=BODY, color=INK, space_after=0, line=110)],
              anchor=MSO_ANCHOR.MIDDLE)
    notes(sl, tidy(SCRIPT[16]))

    # --- 17. where the work goes next -------------------------------------------
    # Closes the talk on reach and next steps, kept separate from the takeaways
    # so the limitations on slide 15 are not what the audience hears last.
    sl = s(17)
    set_title(sl, "Where this goes next")
    drop_body(sl)
    y = lead(sl, [
        "Anywhere a person drives a gripper they cannot feel, this failure is "
        "waiting.",
        "None of these need a new robot, only a better signal on the skin."])
    plate(sl, y, [
        (None, "", "Keyhole surgery",
         "PHOTO: surgeon at a teleoperated console"),
        (None, "", "Glovebox and hot-cell work",
         "PHOTO: remote manipulator handling sealed samples"),
        (None, "", "Disaster and orbital servicing",
         "PHOTO: field or space manipulator on a fragile payload"),
    ], CB - y - 320000)
    notes(sl, tidy(SCRIPT[17]))

    prs.save(OUT)
    print("wrote", OUT)

finally:
    shutil.rmtree(WORK, ignore_errors=True)
